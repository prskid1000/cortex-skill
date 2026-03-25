#!/usr/bin/env python3
"""
Cortex Example: pptx-status-deck
Description: Generate status report PowerPoint from JSON
Tags: powerpoint,pptx,report,status
Captured: 2026-03-25
Source: pptx-status-deck.py

Usage:
  python ~/.claude/skills/cortex/cookbook/pptx-status-deck.py
"""

#!/usr/bin/env python3
"""
Generate a status report PowerPoint deck from JSON data.
Includes title slide, KPI summary, task progress, and next steps.

Usage:
    python pptx-status-deck.py report_data.json [output.pptx]

JSON format:
{
    "title": "Weekly Status Report",
    "subtitle": "Week 13 — March 2026",
    "author": "Team Name",
    "kpis": [
        {"name": "Revenue", "value": "$150K", "change": "+12%", "status": "green"},
        {"name": "Active Users", "value": "5,200", "change": "+8%", "status": "green"},
        {"name": "Open Bugs", "value": "23", "change": "+5", "status": "red"}
    ],
    "tasks": {
        "done": ["Implemented auth flow", "Fixed payment bug"],
        "in_progress": ["Dashboard redesign", "API v2 migration"],
        "blocked": ["Cloud migration — waiting on vendor"]
    },
    "next_steps": ["Complete API v2", "Start load testing", "Review security audit"]
}
"""

import sys
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


STATUS_COLORS = {
    "green": RGBColor(0x28, 0xA7, 0x45),
    "amber": RGBColor(0xFF, 0xA5, 0x00),
    "red": RGBColor(0xDC, 0x35, 0x45),
}


def add_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = data["title"]
    if slide.placeholders[1]:
        slide.placeholders[1].text = f"{data.get('subtitle', '')}\n{data.get('author', '')}"


def add_kpi_slide(prs, kpis):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    slide.shapes.title.text = "Key Metrics"

    cols = len(kpis)
    col_width = Inches(9) // cols
    left_start = Inches(0.5)

    for i, kpi in enumerate(kpis):
        left = left_start + col_width * i
        top = Inches(2)

        # Value
        txBox = slide.shapes.add_textbox(left, top, col_width, Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = kpi["value"]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = STATUS_COLORS.get(kpi.get("status", "green"), RGBColor(0, 0, 0))
        p.alignment = PP_ALIGN.CENTER

        # Change indicator
        txBox2 = slide.shapes.add_textbox(left, Inches(2.8), col_width, Inches(0.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = kpi.get("change", "")
        p2.font.size = Pt(18)
        p2.alignment = PP_ALIGN.CENTER

        # Name
        txBox3 = slide.shapes.add_textbox(left, Inches(3.3), col_width, Inches(0.5))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = kpi["name"]
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(0x70, 0x70, 0x70)
        p3.alignment = PP_ALIGN.CENTER


def add_tasks_slide(prs, tasks):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Task Progress"

    sections = [
        ("Done", tasks.get("done", []), RGBColor(0x28, 0xA7, 0x45)),
        ("In Progress", tasks.get("in_progress", []), RGBColor(0x00, 0x7B, 0xFF)),
        ("Blocked", tasks.get("blocked", []), RGBColor(0xDC, 0x35, 0x45)),
    ]

    for col_idx, (label, items, color) in enumerate(sections):
        left = Inches(0.5 + col_idx * 3.2)
        top = Inches(1.8)

        # Section header
        txBox = slide.shapes.add_textbox(left, top, Inches(2.8), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{label} ({len(items)})"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        # Items
        for i, item in enumerate(items[:6]):
            txBox2 = slide.shapes.add_textbox(left, Inches(2.4 + i * 0.4), Inches(2.8), Inches(0.4))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = f"• {item}"
            p2.font.size = Pt(12)


def add_next_steps_slide(prs, steps):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Next Steps"

    for i, step in enumerate(steps):
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2 + i * 0.5), Inches(8), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{i+1}. {step}"
        p.font.size = Pt(18)


def generate_deck(data, output_path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    add_title_slide(prs, data)
    if data.get("kpis"):
        add_kpi_slide(prs, data["kpis"])
    if data.get("tasks"):
        add_tasks_slide(prs, data["tasks"])
    if data.get("next_steps"):
        add_next_steps_slide(prs, data["next_steps"])

    prs.save(str(output_path))
    print(f"Deck saved: {output_path}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python pptx-status-deck.py report_data.json [output.pptx]")
        sys.exit(1)

    data_path = Path(sys.argv[1])
    with open(data_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    output = sys.argv[2] if len(sys.argv) > 2 else str(data_path.with_suffix(".pptx"))
    generate_deck(data, output)
